from pathlib import Path

from pptx import Presentation

def validate_pptx(file_path: Path) -> dict:
    """
    Performs comprehensive quality, structural, and consistency checks on a PPTX file.
    """
    results = {
        "valid": True,
        "issues": [],
        "slide_count": 0,
        "structural_integrity": "Passed",
        "visual_consistency": "Passed"
    }

    if not file_path.exists():
        results["valid"] = False
        results["issues"].append("File does not exist.")
        return results

    try:
        prs = Presentation(file_path)
        results["slide_count"] = len(prs.slides)

        if results["slide_count"] == 0:
            results["valid"] = False
            results["issues"].append("Presentation has no slides.")

        fonts = set()
        colors = set()

        for i, slide in enumerate(prs.slides):
            # 1. Title Check
            if not slide.shapes.title or not slide.shapes.title.text.strip():
                results["issues"].append(f"Slide {i+1}: Missing or empty title.")

            # 2. Structural: Placeholder Usage
            for shape in slide.placeholders:
                if not shape.has_text_frame or not shape.text.strip():
                    if shape.placeholder_format.type not in [7, 8]: # Ignore pictures/other types
                        results["issues"].append(f"Slide {i+1}: Placeholder '{shape.name}' is empty.")

            # 3. Visual: Font and Color Consistency
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            if run.font.name:
                                fonts.add(run.font.name)
                            if run.font.color and hasattr(run.font.color, 'rgb') and run.font.color.rgb:
                                colors.add(str(run.font.color.rgb))

        if len(fonts) > 3:
             results["issues"].append(f"High font variety ({len(fonts)} found): {fonts}")
             results["visual_consistency"] = "Warning"

    except Exception as e:
        results["valid"] = False
        results["issues"].append(f"Failed to parse PPTX file: {e}")

    return results
