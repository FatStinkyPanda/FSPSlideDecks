from pathlib import Path

from .deck_manager import get_deck_path
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches

def create_enhanced_deck(deck_name: str, slides_content: list) -> Path:
    """
    Generates a PPTX file with support for multiple layouts.
    slides_content: list of dicts like:
    {
        'title': '...',
        'content': '...',
        'layout': 1, # default Title and Content
        'bullet_points': ['...', '...']
    }
    """
    deck_path = get_deck_path(deck_name)
    if not deck_path.exists():
        raise FileNotFoundError(f"Deck folder '{deck_name}' does not exist.")

    prs = Presentation()

    for slide_data in slides_content:
        layout_idx = slide_data.get('layout', 1)
        if layout_idx >= len(prs.slide_layouts):
            layout_idx = 1 # fallback

        slide_layout = prs.slide_layouts[layout_idx]
        slide = prs.slides.add_slide(slide_layout)

        # Handle Title
        if slide.shapes.title:
            slide.shapes.title.text = slide_data.get('title', 'Untitled Slide')

        # Handle Charts
        if 'chart_data' in slide_data:
            c_data = slide_data['chart_data']
            chart_data = CategoryChartData()
            chart_data.categories = c_data.get('categories', [])
            for series_name, values in c_data.get('series', {}).items():
                chart_data.add_series(series_name, values)

            x, y, cx, cy = Inches(0.5), Inches(1.5), Inches(9), Inches(5)
            chart_type = getattr(XL_CHART_TYPE, c_data.get('type', 'COLUMN_CLUSTERED'))
            slide.shapes.add_chart(chart_type, x, y, cx, cy, chart_data)

        # Handle Tables
        elif 'table_data' in slide_data:
            data = slide_data['table_data']
            rows, cols = len(data), len(data[0])
            left, top, width, height = Inches(0.5), Inches(1.5), Inches(9), Inches(5)
            table = slide.shapes.add_table(rows, cols, left, top, width, height).table

            for r in range(rows):
                for c in range(cols):
                    table.cell(r, c).text = str(data[r][c])

        # Handle Content / Body (if not a table)
        elif len(slide.placeholders) > 1:
            body_shape = slide.placeholders[1]
            tf = body_shape.text_frame

            # Simple content text or bullet points
            if 'bullet_points' in slide_data:
                tf.text = slide_data['bullet_points'][0]
                for bp in slide_data['bullet_points'][1:]:
                    p = tf.add_paragraph()
                    p.text = bp
                    p.level = 0
            else:
                tf.text = slide_data.get('content', '')

    output_file = deck_path / "output" / f"{deck_name}.pptx"
    prs.save(output_file)
    return output_file

def add_image_slide(deck_name: str, title_text: str, image_filename: str) -> Path:
    """
    Adds a slide with an image from the assets folder to an existing deck.
    If the deck output doesn't exist, it creates a new one.
    """
    deck_path = get_deck_path(deck_name)
    image_path = deck_path / "assets" / image_filename
    output_file = deck_path / "output" / f"{deck_name}.pptx"

    if not image_path.exists():
        raise FileNotFoundError(f"Asset '{image_filename}' not found in {deck_name}/assets")

    if output_file.exists():
        prs = Presentation(output_file)
    else:
        prs = Presentation()

    # Picture layout (often layout 8 in standard templates)
    # We'll use a blank layout (6) and add shapes manually for more control
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)

    # Add title manually if needed, or use a layout with placeholders
    # Let's keep it simple and just add the picture
    left = top = Inches(1)
    slide.shapes.add_picture(str(image_path), left, top, height=Inches(5))

    # Add title text box
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(1))
    tf = txBox.text_frame
    tf.text = title_text

    prs.save(output_file)
    return output_file
